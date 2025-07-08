"""
PowerPoint Sanitizer
===================

Complete sanitization system for PowerPoint presentations.
Core functionalities:
1. Shape Identification & Extraction
2. AI-Enhanced Sensitive Data Detection  
3. Content Replacement
4. Output: Sanitized PPTX + Analysis Report
"""

import os
import json
import logging
from typing import List, Dict, Any
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx_parser import PPTXParser, SlideData
from openai_analyzer import OpenAIAnalyzer, Detection

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class SanitizationReport:
    """Report of sanitization results."""
    original_file: str
    sanitized_file: str
    total_slides: int
    total_detections: int
    detections_by_slide: Dict[int, List[Detection]] = field(default_factory=dict)
    categories_summary: Dict[str, int] = field(default_factory=dict)


class PowerPointSanitizer:
    """Complete PowerPoint sanitization system."""
    
    def __init__(self, openai_api_key: str, images_dir: str = "data/pngs"):
        """Initialize sanitizer with OpenAI API key and images directory."""
        self.parser = PPTXParser()
        self.analyzer = OpenAIAnalyzer(api_key=openai_api_key)
        self.images_dir = Path(images_dir)
        self.logger = logging.getLogger(__name__)
        
    def sanitize_presentation(self, input_file: str, output_file: str = None) -> SanitizationReport:
        """
        Complete sanitization workflow.
        
        Args:
            input_file: Path to input PowerPoint file
            output_file: Path for sanitized output file (optional)
            
        Returns:
            SanitizationReport with results
        """
        # Set default output file name
        if not output_file:
            input_path = Path(input_file)
            output_file = str(input_path.parent / f"{input_path.stem}_sanitized{input_path.suffix}")
        
        self.logger.info(f"Starting sanitization of {input_file}")
        
        # 1. Shape Identification & Extraction
        slides_data = self.parser.parse_presentation(input_file)
        self.logger.info(f"Extracted data from {len(slides_data)} slides")
        
        # 2. AI-Enhanced Sensitive Data Detection
        all_detections = {}
        for slide in slides_data:
            slide_image_path = self.images_dir / f"slide_{slide.slide_number:02d}.png"
            
            if slide_image_path.exists():
                detections = self._analyze_slide(slide, slide_image_path)
                all_detections[slide.slide_number] = detections
                self.logger.info(f"Slide {slide.slide_number}: {len(detections.detections)} detections")
            else:
                self.logger.warning(f"Image not found for slide {slide.slide_number}: {slide_image_path}")
                all_detections[slide.slide_number] = []
        
        # 3. Content Replacement
        self._replace_content_in_file(input_file, output_file, all_detections)
        
        # 4. Generate Report
        report = self._generate_report(input_file, output_file, slides_data, all_detections)
        self._save_report(report, output_file)
        
        self.logger.info(f"Sanitization completed. Output: {output_file}")
        return report
    
    def _analyze_slide(self, slide_data: SlideData, image_path: Path):
        """Analyze a single slide for sensitive content."""
        if not slide_data.text_content:
            return []
            
        try:
            detections = self.analyzer.analyze_slide(
                slide_text=slide_data.text_content,
                image_path=str(image_path)
            )
            return detections
        except Exception as e:
            self.logger.error(f"Error analyzing slide {slide_data.slide_number}: {e}")
            return []
    
    def _replace_content_in_file(self, input_file: str, output_file: str, all_detections: Dict[int, Any]):
        """Replace sensitive content in PowerPoint file."""
        try:
            # Load presentation
            presentation = Presentation(input_file)
            
            # Process each slide
            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                
                if slide_number in all_detections:
                    detections = all_detections[slide_number]
                    if hasattr(detections, 'detections'):
                        self._replace_slide_content(slide, detections.detections)
                    else:
                        self._replace_slide_content(slide, detections)
            
            # Save sanitized presentation
            presentation.save(output_file)
            self.logger.info(f"Saved sanitized presentation to {output_file}")
            
        except Exception as e:
            self.logger.error(f"Error replacing content: {e}")
            raise
    
    def _replace_slide_content(self, slide, detections: List[Detection]):
        """Replace content in a single slide."""
        if not detections:
            return
            
        # Create replacement map - sort by length (longest first) to avoid partial replacements
        replacement_map = {}
        for detection in detections:
            if hasattr(detection, 'original') and hasattr(detection, 'replacement'):
                replacement_map[detection.original] = detection.replacement
            elif hasattr(detection, 'text') and hasattr(detection, 'replacement'):
                replacement_map[detection.text] = detection.replacement
        
        if not replacement_map:
            return
        
        # Sort by length (longest first) to prevent partial replacements
        sorted_replacements = sorted(replacement_map.items(), key=lambda x: len(x[0]), reverse=True)
        
        self.logger.info(f"Applying {len(sorted_replacements)} replacements to slide")
        for original, replacement in sorted_replacements:
            self.logger.debug(f"Will replace: '{original}' -> '{replacement}'")
        
        # Process all shapes in the slide
        for shape in slide.shapes:
            self._replace_shape_text(shape, sorted_replacements)
        
        # Additional pass: handle any text that might be in combined elements
        self._replace_slide_text_comprehensive(slide, sorted_replacements)
    
    def _replace_shape_text(self, shape, sorted_replacements: List[tuple]):
        """Replace text in a single shape."""
        try:
            # Handle text frames
            if shape.has_text_frame:
                self._replace_text_frame(shape.text_frame, sorted_replacements)
            
            # Handle tables
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self._replace_table_text(shape.table, sorted_replacements)
                
            # Handle charts (text in chart elements)
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                # Charts may have text in titles, labels, etc.
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    self._replace_text_frame(shape.text_frame, sorted_replacements)
                    
        except Exception as e:
            self.logger.warning(f"Error replacing text in shape: {e}")
    
    def _replace_text_frame(self, text_frame, sorted_replacements: List[tuple]):
        """Replace text in a text frame."""
        if not text_frame or not text_frame.paragraphs:
            return
            
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text:
                    continue
                    
                original_text = run.text
                new_text = original_text
                
                # Apply all replacements in order (longest first)
                for original, replacement in sorted_replacements:
                    if original and original in new_text:
                        new_text = new_text.replace(original, replacement)
                        self.logger.debug(f"In run, replaced: '{original}' with '{replacement}'")
                
                # Update the run text if changes were made
                if new_text != original_text:
                    run.text = new_text
                    self.logger.info(f"Updated run: '{original_text}' -> '{new_text}'")
    
    def _replace_slide_text_comprehensive(self, slide, sorted_replacements: List[tuple]):
        """Additional comprehensive text replacement for the entire slide."""
        try:
            # Get all text content from the slide
            all_text_shapes = []
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text_shapes.append(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # Add table cells
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                all_text_shapes.append(cell)
            
            # Apply replacements to each text shape
            for text_shape in all_text_shapes:
                if hasattr(text_shape, 'text_frame'):
                    text_frame = text_shape.text_frame
                elif hasattr(text_shape, 'text'):
                    # For table cells
                    continue  # Already handled in _replace_table_text
                else:
                    continue
                
                # Get the full text content
                full_text = text_frame.text
                if not full_text:
                    continue
                
                # Apply replacements
                new_full_text = full_text
                replacements_made = []
                
                for original, replacement in sorted_replacements:
                    if original and original in new_full_text:
                        new_full_text = new_full_text.replace(original, replacement)
                        replacements_made.append((original, replacement))
                
                # If replacements were made, update the text frame
                if replacements_made and new_full_text != full_text:
                    # Clear existing text and set new text
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = new_full_text
                    
                    self.logger.info(f"Comprehensive replacement made: {len(replacements_made)} items")
                    for orig, repl in replacements_made:
                        self.logger.debug(f"  '{orig}' -> '{repl}'")
                        
        except Exception as e:
            self.logger.warning(f"Error in comprehensive text replacement: {e}")
    
    def _replace_table_text(self, table, sorted_replacements: List[tuple]):
        """Replace text in table cells."""
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    self._replace_text_frame(cell.text_frame, sorted_replacements)
    
    def _generate_report(self, input_file: str, output_file: str, 
                        slides_data: List[SlideData], all_detections: Dict[int, Any]) -> SanitizationReport:
        """Generate sanitization report."""
        
        # Count total detections and categorize
        total_detections = 0
        categories_summary = {}
        detections_by_slide = {}
        
        for slide_number, detections in all_detections.items():
            if hasattr(detections, 'detections'):
                detection_list = detections.detections
            else:
                detection_list = detections if isinstance(detections, list) else []
            
            detections_by_slide[slide_number] = detection_list
            total_detections += len(detection_list)
            
            # Count categories
            for detection in detection_list:
                category = getattr(detection, 'category', 'unknown')
                categories_summary[category] = categories_summary.get(category, 0) + 1
        
        return SanitizationReport(
            original_file=input_file,
            sanitized_file=output_file,
            total_slides=len(slides_data),
            total_detections=total_detections,
            detections_by_slide=detections_by_slide,
            categories_summary=categories_summary
        )
    
    def _save_report(self, report: SanitizationReport, output_file: str):
        """Save sanitization report as JSON."""
        report_file = Path(output_file).with_suffix('.json')
        
        # Convert report to serializable format
        report_data = {
            "original_file": report.original_file,
            "sanitized_file": report.sanitized_file,
            "total_slides": report.total_slides,
            "total_detections": report.total_detections,
            "categories_summary": report.categories_summary,
            "detections_by_slide": {
                str(slide_num): [
                    {
                        "original": getattr(d, 'original', getattr(d, 'text', '')),
                        "replacement": getattr(d, 'replacement', ''),
                        "category": getattr(d, 'category', 'unknown'),
                        "reason": getattr(d, 'reason', '')
                    } for d in detections
                ] for slide_num, detections in report.detections_by_slide.items()
            }
        }
        
        with open(report_file, 'w', encoding='utf-8') as f:
            json.dump(report_data, f, indent=2, ensure_ascii=False)
        
        self.logger.info(f"Saved sanitization report to {report_file}")
    
    def print_summary(self, report: SanitizationReport):
        """Print sanitization summary."""
        print(f"\n=== SANITIZATION SUMMARY ===")
        print(f"Original file: {report.original_file}")
        print(f"Sanitized file: {report.sanitized_file}")
        print(f"Total slides: {report.total_slides}")
        print(f"Total detections: {report.total_detections}")
        
        print(f"\nDetections by category:")
        for category, count in report.categories_summary.items():
            print(f"  {category}: {count}")
        
        print(f"\nDetections by slide:")
        for slide_num, detections in report.detections_by_slide.items():
            if detections:
                print(f"  Slide {slide_num}: {len(detections)} detections")


def main():
    """Main sanitization workflow."""
    
    # Configuration
    INPUT_FILE = "data/Take-home.pptx"
    OUTPUT_FILE = "data/Take-home_sanitized.pptx"
    OPENAI_API_KEY = "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"  # Replace with your actual API key
    IMAGES_DIR = "data/pngs"
    
    try:
        # Initialize sanitizer
        sanitizer = PowerPointSanitizer(
            openai_api_key=OPENAI_API_KEY,
            images_dir=IMAGES_DIR
        )
        
        # Run sanitization
        report = sanitizer.sanitize_presentation(INPUT_FILE, OUTPUT_FILE)
        
        # Print summary
        sanitizer.print_summary(report)
        
        print(f"\nSanitization completed successfully!")
        print(f"Sanitized file: {report.sanitized_file}")
        print(f"Report file: {Path(report.sanitized_file).with_suffix('.json')}")
        
    except Exception as e:
        logger.error(f"Sanitization failed: {e}")
        raise


if __name__ == "__main__":
    main()
