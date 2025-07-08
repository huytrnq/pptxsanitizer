"""
PPTX Parser MVP
==============

Simple class to extract content from PowerPoint files.
"""

import logging
from typing import List, Dict, Any
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class SlideData:
    """Data from a single slide."""

    slide_number: int
    title: str = ""
    text_content: List[str] = field(default_factory=list)
    images_count: int = 0
    charts_count: int = 0
    tables_count: int = 0


class PPTXParser:
    """Simple PowerPoint parser."""

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def parse_presentation(self, file_path: str) -> List[SlideData]:
        """Parse a PowerPoint file and return slide data."""
        try:
            # Load presentation
            presentation = Presentation(file_path)
            self.logger.info(
                f"Loaded presentation with {len(presentation.slides)} slides"
            )

            slides_data = []

            # Process each slide
            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                slide_data = self._parse_slide(slide, slide_number)
                slides_data.append(slide_data)

                self.logger.info(
                    f"Parsed slide {slide_number}: {len(slide_data.text_content)} text elements"
                )

            return slides_data

        except Exception as e:
            self.logger.error(f"Failed to parse presentation: {e}")
            raise

    def _parse_slide(self, slide, slide_number: int) -> SlideData:
        """Parse a single slide."""
        slide_data = SlideData(slide_number=slide_number)

        # Get slide title
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_data.title = slide.shapes.title.text.strip()
        except:
            slide_data.title = f"Slide {slide_number}"

        # Process all shapes
        for shape in slide.shapes:
            self._process_shape(shape, slide_data)

        return slide_data

    def _process_shape(self, shape, slide_data: SlideData):
        """Process a single shape from the slide."""
        try:
            # Text content
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_data.text_content.append(shape.text_frame.text.strip())

            # Count different shape types
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data.images_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                slide_data.charts_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                slide_data.tables_count += 1
                # Extract table text
                self._extract_table_text(shape, slide_data)

        except Exception as e:
            self.logger.warning(f"Error processing shape: {e}")

    def _extract_table_text(self, table_shape, slide_data: SlideData):
        """Extract text from table cells."""
        try:
            table = table_shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame and cell.text_frame.text.strip():
                        slide_data.text_content.append(cell.text_frame.text.strip())
        except Exception as e:
            self.logger.warning(f"Error extracting table text: {e}")


def main():
    """Test the parser."""
    import sys

    file_path = "data/test.pptx"

    try:
        parser = PPTXParser()
        slides = parser.parse_presentation(file_path)

        print(f"\n=== PRESENTATION SUMMARY ===")
        print(f"Total slides: {len(slides)}")

        for slide in slides:
            print(f"\nSlide {slide.slide_number}: {slide.title}")
            print(f"  Text elements: {len(slide.text_content)}")
            print(f"  Images: {slide.images_count}")
            print(f"  Charts: {slide.charts_count}")
            print(f"  Tables: {slide.tables_count}")

            # Show first few text elements
            for i, text in enumerate(slide.text_content):
                preview = text[:50] + "..." if len(text) > 50 else text
                print(f"    Text {i+1}: {preview}")

    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    main()
