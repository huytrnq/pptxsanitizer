"""
PPTX Processor
==============

Comprehensive class for PowerPoint file processing including:
- Content extraction
- Text replacement/sanitization
- File manipulation
"""

import logging
import re
from typing import List, Dict, Any, Tuple
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class SlideData:
    """Data from a single slide."""

    slide_number: int
    title: str = ""
    text_content: List[str] = field(default_factory=list)
    images_count: int = 0
    charts_count: int = 0
    tables_count: int = 0


@dataclass
class Detection:
    """Detection object for compatibility."""

    original: str
    replacement: str
    category: str = ""
    reason: str = ""


class PPTXProcessor:
    """Comprehensive PowerPoint processor with extraction and replacement capabilities."""

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def parse_presentation(self, file_path: str) -> List[SlideData]:
        """Parse a PowerPoint file and return slide data."""
        try:
            # Load presentation
            presentation = Presentation(file_path)
            self.logger.info(
                f"Loaded presentation with {len(presentation.slides)} slides"
            )

            slides_data = []

            # Process each slide
            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                slide_data = self._parse_slide(slide, slide_number)
                slides_data.append(slide_data)

                self.logger.info(
                    f"Parsed slide {slide_number}: {len(slide_data.text_content)} text elements"
                )

            return slides_data

        except Exception as e:
            self.logger.error(f"Failed to parse presentation: {e}")
            raise

    def _parse_slide(self, slide, slide_number: int) -> SlideData:
        """Parse a single slide."""
        slide_data = SlideData(slide_number=slide_number)

        # Get slide title
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_data.title = slide.shapes.title.text.strip()
        except:
            slide_data.title = f"Slide {slide_number}"

        # Process all shapes
        for shape in slide.shapes:
            self._process_shape(shape, slide_data)

        return slide_data

    def _process_shape(self, shape, slide_data: SlideData):
        """Process a single shape from the slide."""
        try:
            # Text content
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_data.text_content.append(shape.text_frame.text.strip())

            # Count different shape types
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data.images_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                slide_data.charts_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                slide_data.tables_count += 1
                # Extract table text
                self._extract_table_text(shape, slide_data)

        except Exception as e:
            self.logger.warning(f"Error processing shape: {e}")

    def _extract_table_text(self, table_shape, slide_data: SlideData):
        """Extract text from table cells."""
        try:
            table = table_shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame and cell.text_frame.text.strip():
                        slide_data.text_content.append(cell.text_frame.text.strip())
        except Exception as e:
            self.logger.warning(f"Error extracting table text: {e}")

    def apply_replacements_to_file(
        self, input_file: str, output_file: str, all_detections: Dict[int, List]
    ) -> Dict[str, Any]:
        """Apply text replacements to PowerPoint file and save sanitized version."""
        try:
            from pptx import Presentation

            # Load presentation
            presentation = Presentation(input_file)
            self.logger.info(f"Loaded presentation for replacement: {input_file}")

            # Process each slide
            total_replacements = 0
            replacements_by_slide = {}

            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1

                if slide_number in all_detections:
                    detections = all_detections[slide_number]
                    # Use your existing _apply_replacements_to_slide method
                    replacements_made = self._apply_replacements_to_slide(
                        slide, detections
                    )
                    total_replacements += replacements_made
                    replacements_by_slide[slide_number] = replacements_made
                    self.logger.info(
                        f"Slide {slide_number}: {replacements_made} replacements applied"
                    )
                else:
                    replacements_by_slide[slide_number] = 0
                    self.logger.info(f"Slide {slide_number}: No detections to apply")

            # Save sanitized presentation
            presentation.save(output_file)
            self.logger.info(f"Saved sanitized presentation to {output_file}")
            self.logger.info(f"Total replacements applied: {total_replacements}")

            # Return detailed results instead of just True/False
            return {
                "success": True,
                "total_replacements": total_replacements,
                "replacements_by_slide": replacements_by_slide,
            }

        except Exception as e:
            self.logger.error(f"Error applying replacements: {e}")
            return {
                "success": False,
                "total_replacements": 0,
                "replacements_by_slide": {},
                "error": str(e),
            }

    def _apply_replacements_to_slide(self, slide, detections: List[Detection]) -> int:
        """Apply replacements to a single slide."""
        if not detections:
            return 0

        # Create sorted replacement list (longest first to avoid partial replacements)
        replacements = []
        for detection in detections:
            if hasattr(detection, "original") and hasattr(detection, "replacement"):
                replacements.append((detection.original, detection.replacement))
            elif hasattr(detection, "text") and hasattr(detection, "replacement"):
                replacements.append((detection.text, detection.replacement))

        if not replacements:
            return 0

        # Sort by length (longest first)
        sorted_replacements = sorted(
            replacements, key=lambda x: len(x[0]), reverse=True
        )

        self.logger.info(f"Applying {len(sorted_replacements)} replacements:")
        for original, replacement in sorted_replacements:
            self.logger.info(f"  '{original}' -> '{replacement}'")

        # Debug: Show current slide content
        self._debug_slide_content(slide)

        # Apply replacements to all shapes
        total_replacements = 0
        for shape in slide.shapes:
            replacements_made = self._apply_replacements_to_shape(
                shape, sorted_replacements
            )
            total_replacements += replacements_made

        return total_replacements

    def _apply_replacements_to_shape(
        self, shape, sorted_replacements: List[Tuple[str, str]]
    ) -> int:
        """Apply replacements to a single shape."""
        try:
            replacements_made = 0

            # Handle text frames
            if shape.has_text_frame:
                replacements_made += self._apply_replacements_to_text_frame(
                    shape.text_frame, sorted_replacements
                )

            # Handle tables
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                replacements_made += self._apply_replacements_to_table(
                    shape.table, sorted_replacements
                )

            # Handle charts (text in chart elements)
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    replacements_made += self._apply_replacements_to_text_frame(
                        shape.text_frame, sorted_replacements
                    )

            return replacements_made

        except Exception as e:
            self.logger.warning(f"Error applying replacements to shape: {e}")
            return 0

    def _apply_replacements_to_text_frame(
        self, text_frame, sorted_replacements: List[Tuple[str, str]]
    ) -> int:
        """Apply replacements to a text frame."""
        if not text_frame or not text_frame.paragraphs:
            return 0

        replacements_made = 0

        # Get full text content first
        full_text = text_frame.text
        if not full_text:
            return 0

        self.logger.debug(f"Processing text frame: '{full_text}'")

        # Apply replacements to full text
        new_full_text = full_text
        applied_replacements = []

        for original, replacement in sorted_replacements:
            if original and original in new_full_text:
                new_full_text = new_full_text.replace(original, replacement)
                applied_replacements.append((original, replacement))
                self.logger.info(f"  Replaced: '{original}' -> '{replacement}'")

        # If no exact matches, try fuzzy matching
        if not applied_replacements:
            fuzzy_results = self._apply_fuzzy_replacements(
                full_text, sorted_replacements
            )
            new_full_text = fuzzy_results["new_text"]
            applied_replacements = fuzzy_results["replacements"]

        # Update text frame if changes were made
        if applied_replacements and new_full_text != full_text:
            try:
                # Clear and rewrite the entire text frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = new_full_text

                replacements_made = len(applied_replacements)
                self.logger.info(
                    f"  Updated text frame: {replacements_made} replacements"
                )
                self.logger.debug(f"    Original: '{full_text}'")
                self.logger.debug(f"    New: '{new_full_text}'")

            except Exception as e:
                self.logger.error(f"Error updating text frame: {e}")
                return 0

        return replacements_made

    def _apply_replacements_to_table(
        self, table, sorted_replacements: List[Tuple[str, str]]
    ) -> int:
        """Apply replacements to table cells."""
        replacements_made = 0

        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    replacements_made += self._apply_replacements_to_text_frame(
                        cell.text_frame, sorted_replacements
                    )

        return replacements_made

    def _apply_fuzzy_replacements(
        self, text: str, sorted_replacements: List[Tuple[str, str]]
    ) -> Dict[str, Any]:
        """Apply fuzzy/flexible text replacements."""
        new_text = text
        applied_replacements = []

        self.logger.info(f"Attempting fuzzy matching for: '{text}'")

        for original, replacement in sorted_replacements:
            if not original:
                continue

            # Normalize both texts
            normalized_text = self._normalize_text_for_matching(new_text)
            normalized_original = self._normalize_text_for_matching(original)

            # Try different matching strategies
            if normalized_original in normalized_text:
                # Simple normalized match
                new_text = new_text.replace(original, replacement)
                applied_replacements.append((original, replacement))
                self.logger.info(
                    f"  Fuzzy match (normalized): '{original}' -> '{replacement}'"
                )

            elif self._is_flexible_match(normalized_text, normalized_original):
                # Flexible regex match
                pattern = self._create_flexible_pattern(normalized_original)
                new_text = re.sub(pattern, replacement, new_text, flags=re.IGNORECASE)
                applied_replacements.append((original, replacement))
                self.logger.info(
                    f"  Fuzzy match (flexible): '{original}' -> '{replacement}'"
                )

        return {"new_text": new_text, "replacements": applied_replacements}

    def _normalize_text_for_matching(self, text: str) -> str:
        """Normalize text for better matching."""
        if not text:
            return ""

        # Remove extra whitespace
        text = re.sub(r"\s+", " ", text.strip())
        # Handle special characters
        text = text.replace("\u00a0", " ")  # Non-breaking space
        text = text.replace("\u2013", "-")  # En dash
        text = text.replace("\u2014", "-")  # Em dash
        text = text.replace("\u201c", '"')  # Left double quote
        text = text.replace("\u201d", '"')  # Right double quote
        text = text.replace("\u2018", "'")  # Left single quote
        text = text.replace("\u2019", "'")  # Right single quote

        return text

    def _is_flexible_match(self, text: str, search_term: str) -> bool:
        """Check if search term matches with flexible whitespace/formatting."""
        # Create flexible pattern
        pattern = self._create_flexible_pattern(search_term)
        return bool(re.search(pattern, text, re.IGNORECASE))

    def _create_flexible_pattern(self, search_term: str) -> str:
        """Create a flexible regex pattern for matching."""
        # Escape special regex characters
        escaped = re.escape(search_term)
        # Replace escaped spaces with flexible whitespace pattern
        flexible = escaped.replace(r"\ ", r"\s*")
        return flexible

    def _debug_slide_content(self, slide):
        """Debug method to show slide content."""
        self.logger.info("Current slide content:")

        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                self.logger.info(
                    f"  Shape {shape_idx} (text): '{shape.text_frame.text.strip()}'"
                )
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self.logger.info(f"  Shape {shape_idx} (table):")
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        if cell.text_frame and cell.text_frame.text.strip():
                            self.logger.info(
                                f"    Cell[{row_idx},{cell_idx}]: '{cell.text_frame.text.strip()}'"
                            )
            else:
                self.logger.info(f"  Shape {shape_idx} ({shape.shape_type}): No text")

    def debug_presentation_structure(self, file_path: str):
        """Debug method to examine presentation structure."""
        try:
            presentation = Presentation(file_path)

            print(f"\n=== PRESENTATION STRUCTURE DEBUG ===")
            print(f"File: {file_path}")
            print(f"Total slides: {len(presentation.slides)}")

            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                print(f"\n--- Slide {slide_number} ---")

                for shape_idx, shape in enumerate(slide.shapes):
                    print(f"Shape {shape_idx}:")
                    print(f"  Type: {shape.shape_type}")
                    print(f"  Has text frame: {shape.has_text_frame}")

                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        print(f"  Text: '{text}'")
                        print(f"  Paragraphs: {len(shape.text_frame.paragraphs)}")

                        for para_idx, paragraph in enumerate(
                            shape.text_frame.paragraphs
                        ):
                            print(
                                f"    Para {para_idx}: '{paragraph.text}' (runs: {len(paragraph.runs)})"
                            )

                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = shape.table
                        print(f"  Table: {len(table.rows)}x{len(table.columns)}")

                        for row_idx, row in enumerate(table.rows):
                            for cell_idx, cell in enumerate(row.cells):
                                if cell.text_frame and cell.text_frame.text.strip():
                                    print(
                                        f"    Cell[{row_idx},{cell_idx}]: '{cell.text_frame.text}'"
                                    )

        except Exception as e:
            print(f"Error debugging presentation: {e}")


def main():
    """Test the processor."""
    file_path = "data/test.pptx"

    try:
        processor = PPTXProcessor()

        # Test parsing
        slides = processor.parse_presentation(file_path)
        print(f"\n=== PRESENTATION SUMMARY ===")
        print(f"Total slides: {len(slides)}")

        for slide in slides:
            print(f"\nSlide {slide.slide_number}: {slide.title}")
            print(f"  Text elements: {len(slide.text_content)}")
            print(f"  Images: {slide.images_count}")
            print(f"  Charts: {slide.charts_count}")
            print(f"  Tables: {slide.tables_count}")

            # Show text content
            for i, text in enumerate(slide.text_content):
                preview = text[:50] + "..." if len(text) > 50 else text
                print(f"    Text {i+1}: {preview}")

        # Test debug method
        processor.debug_presentation_structure(file_path)

    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    main()
