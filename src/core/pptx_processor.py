"""
PPTX Processor
==============

Comprehensive class for PowerPoint file processing including:
- Content extraction from slides, text frames, tables, and charts
- Text replacement and sanitization with formatting preservation
- File manipulation and sanitized output generation
"""

import logging
from typing import List, Dict, Any, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models.slide_data import SlideData
from ..models.detection import Detection
from ..utils.text_processing import TextProcessor


class PPTXProcessor:
    """
    Comprehensive PowerPoint processor with extraction and replacement capabilities.
    
    This class provides methods to:
    - Parse PowerPoint presentations and extract structured data
    - Apply text replacements while preserving formatting
    - Handle various PowerPoint elements (text frames, tables, charts)
    - Generate sanitized versions of presentations
    
    Attributes:
        logger: Logger instance for tracking operations
        text_processor: TextProcessor instance for advanced text operations
    """

    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.text_processor = TextProcessor()

    def parse_presentation(self, file_path: str) -> List[SlideData]:
        """
        Parse a PowerPoint file and return structured slide data.
        
        Extracts content from all slides including text, images, charts, and tables.
        Each slide is processed to create a SlideData object containing all relevant
        information for analysis and sanitization.
        
        Args:
            file_path (str): Path to the PowerPoint file to parse
            
        Returns:
            List[SlideData]: List of SlideData objects, one for each slide
            
        Raises:
            Exception: If the presentation cannot be loaded or parsed
        """
        try:
            # Load presentation
            presentation = Presentation(file_path)
            self.logger.info(
                f"Loaded presentation with {len(presentation.slides)} slides"
            )

            slides_data = []

            # Process each slide
            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1
                slide_data = self._parse_slide(slide, slide_number)
                slides_data.append(slide_data)

                self.logger.info(
                    f"Parsed slide {slide_number}: {len(slide_data.text_content)} text elements"
                )

            return slides_data

        except Exception as e:
            self.logger.error(f"Failed to parse presentation: {e}")
            raise

    def _parse_slide(self, slide, slide_number: int) -> SlideData:
        """
        Parse a single slide and extract all relevant content.
        
        Processes all shapes on the slide to extract text content, count elements,
        and gather metadata. Creates a comprehensive SlideData object.
        
        Args:
            slide: The python-pptx slide object to parse
            slide_number (int): The 1-based slide number for identification
            
        Returns:
            SlideData: Object containing all extracted slide information
        """
        slide_data = SlideData(slide_number=slide_number)

        # Get slide title
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_data.title = slide.shapes.title.text.strip()
        except:
            slide_data.title = f"Slide {slide_number}"

        # Process all shapes
        for shape in slide.shapes:
            self._process_shape(shape, slide_data)

        return slide_data

    def _process_shape(self, shape, slide_data: SlideData):
        """
        Process a single shape from the slide and extract relevant data.
        
        Handles different types of shapes including text frames, images, charts,
        and tables. Extracts text content and updates element counts in slide_data.
        
        Args:
            shape: The python-pptx shape object to process
            slide_data (SlideData): The slide data object to update with extracted information
        """
        try:
            # Text content
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_data.text_content.append(shape.text_frame.text.strip())

            # Count different shape types
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data.images_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                slide_data.charts_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                slide_data.tables_count += 1
                # Extract table text
                self._extract_table_text(shape, slide_data)

        except Exception as e:
            self.logger.warning(f"Error processing shape: {e}")

    def _extract_table_text(self, table_shape, slide_data: SlideData):
        """
        Extract text from table cells and add to slide data.
        
        Iterates through all table rows and cells to extract text content,
        adding non-empty text to the slide's text_content list.
        
        Args:
            table_shape: The python-pptx table shape object
            slide_data (SlideData): The slide data object to update with table text
        """
        try:
            table = table_shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame and cell.text_frame.text.strip():
                        slide_data.text_content.append(cell.text_frame.text.strip())
        except Exception as e:
            self.logger.warning(f"Error extracting table text: {e}")

    def apply_replacements_to_file(
        self, input_file: str, output_file: str, all_detections: Dict[int, List]
    ) -> Dict[str, Any]:
        """
        Apply text replacements to PowerPoint file and save sanitized version.
        
        Loads the input presentation, applies all specified text replacements
        while preserving formatting, and saves the result to the output file.
        
        Args:
            input_file (str): Path to the input PowerPoint file
            output_file (str): Path where the sanitized file will be saved
            all_detections (Dict[int, List]): Dictionary mapping slide numbers to 
                lists of Detection objects containing replacement information
                
        Returns:
            Dict[str, Any]: Result dictionary containing:
                - success (bool): Whether the operation succeeded
                - total_replacements (int): Total number of replacements made
                - replacements_by_slide (Dict[int, int]): Replacements per slide
                - error (str, optional): Error message if operation failed
        """
        try:
            # Load presentation
            presentation = Presentation(input_file)
            self.logger.info(f"Loaded presentation for replacement: {input_file}")

            # Process each slide
            total_replacements = 0
            replacements_by_slide = {}

            for slide_idx, slide in enumerate(presentation.slides):
                slide_number = slide_idx + 1

                if slide_number in all_detections:
                    detections = all_detections[slide_number]
                    replacements_made = self._apply_replacements_to_slide(
                        slide, detections
                    )
                    total_replacements += replacements_made
                    replacements_by_slide[slide_number] = replacements_made
                    self.logger.info(
                        f"Slide {slide_number}: {replacements_made} replacements applied"
                    )
                else:
                    replacements_by_slide[slide_number] = 0
                    self.logger.info(f"Slide {slide_number}: No detections to apply")

            # Save sanitized presentation
            presentation.save(output_file)
            self.logger.info(f"Saved sanitized presentation to {output_file}")
            self.logger.info(f"Total replacements applied: {total_replacements}")

            return {
                "success": True,
                "total_replacements": total_replacements,
                "replacements_by_slide": replacements_by_slide,
            }

        except Exception as e:
            self.logger.error(f"Error applying replacements: {e}")
            return {
                "success": False,
                "total_replacements": 0,
                "replacements_by_slide": {},
                "error": str(e),
            }

    def _apply_replacements_to_slide(self, slide, detections: List[Detection]) -> int:
        """
        Apply replacements to a single slide.
        
        Processes all detection objects for the slide, extracts replacement pairs,
        and applies them to all shapes on the slide. Replacements are sorted by
        length (longest first) to avoid partial replacement issues.
        
        Args:
            slide: The python-pptx slide object to modify
            detections (List[Detection]): List of detection objects containing
                original text and replacement text
                
        Returns:
            int: Total number of replacements made on this slide
        """
        if not detections:
            return 0

        # Create sorted replacement list (longest first to avoid partial replacements)
        replacements = []
        for detection in detections:
            if hasattr(detection, "original") and hasattr(detection, "replacement"):
                replacements.append((detection.original, detection.replacement))
            elif hasattr(detection, "text") and hasattr(detection, "replacement"):
                replacements.append((detection.text, detection.replacement))

        if not replacements:
            return 0

        # Sort by length (longest first)
        sorted_replacements = sorted(
            replacements, key=lambda x: len(x[0]), reverse=True
        )

        self.logger.info(f"Applying {len(sorted_replacements)} replacements:")
        for original, replacement in sorted_replacements:
            self.logger.info(f"  '{original}' -> '{replacement}'")

        # Apply replacements to all shapes
        total_replacements = 0
        for shape in slide.shapes:
            replacements_made = self._apply_replacements_to_shape(
                shape, sorted_replacements
            )
            total_replacements += replacements_made

        return total_replacements
    
    
    def _apply_replacements_to_shape(
        self, shape, sorted_replacements: List[Tuple[str, str]]
    ) -> int:
        """
        Apply replacements to a single shape.
        
        Handles different shape types (text frames, tables, charts) and applies
        the specified text replacements while preserving formatting.
        
        Args:
            shape: The python-pptx shape object to modify
            sorted_replacements (List[Tuple[str, str]]): List of (original, replacement)
                tuples sorted by length (longest first)
                
        Returns:
            int: Number of replacements made in this shape
        """
        try:
            replacements_made = 0

            # Handle text frames
            if shape.has_text_frame:
                replacements_made += self._apply_replacements_to_text_frame(
                    shape.text_frame, sorted_replacements
                )

            # Handle tables
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                replacements_made += self._apply_replacements_to_table(
                    shape.table, sorted_replacements
                )

            # Handle charts (text in chart elements)
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    replacements_made += self._apply_replacements_to_text_frame(
                        shape.text_frame, sorted_replacements
                    )

            return replacements_made

        except Exception as e:
            self.logger.warning(f"Error applying replacements to shape: {e}")
            return 0

    def _apply_replacements_to_text_frame(
        self, text_frame, sorted_replacements: List[Tuple[str, str]]
    ) -> int:
        """
        Apply replacements to a text frame while preserving formatting.
        
        Uses a two-stage approach:
        1. Try run-by-run replacement to preserve detailed formatting
        2. Fall back to full text replacement with basic formatting preservation
        
        The method handles complex formatting scenarios including font properties,
        colors, and paragraph structure. It also supports fuzzy matching when
        exact replacements are not found.
        
        Args:
            text_frame: The python-pptx text frame object to modify
            sorted_replacements (List[Tuple[str, str]]): List of (original, replacement)
                tuples sorted by length (longest first)
                
        Returns:
            int: Number of replacements made in this text frame
        """
        if not text_frame or not text_frame.paragraphs:
            return 0

        replacements_made = 0

        # Try run-by-run replacement first to preserve formatting
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text:
                    continue
                
                original_text = run.text
                new_text = original_text
                
                # Apply replacements
                for original, replacement in sorted_replacements:
                    if original and original in new_text:
                        new_text = new_text.replace(original, replacement)
                        self.logger.info(f"  Replaced: '{original}' -> '{replacement}'")
                
                # Update text while preserving formatting
                if new_text != original_text:
                    # Store original formatting properties
                    original_font_size = run.font.size
                    original_font_name = run.font.name
                    original_bold = run.font.bold
                    original_italic = run.font.italic
                    # Safe color handling
                    original_color = None
                    try:
                        if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                            original_color = run.font.color.rgb
                    except:
                        original_color = None
                    
                    # Update the text
                    run.text = new_text
                    
                    # Restore original formatting
                    if original_font_size:
                        run.font.size = original_font_size
                    if original_font_name:
                        run.font.name = original_font_name
                    if original_bold is not None:
                        run.font.bold = original_bold
                    if original_italic is not None:
                        run.font.italic = original_italic
                    if original_color:
                        try:
                            run.font.color.rgb = original_color
                        except:
                            pass  # Skip color restoration if it fails
                    
                    replacements_made += 1

        # If no run-level replacements worked, fall back to full text replacement
        if not replacements_made:
            # Get full text content
            full_text = text_frame.text
            if not full_text:
                return 0

            self.logger.debug(f"Processing text frame: '{full_text}'")

            # Apply replacements to full text
            new_full_text = full_text
            applied_replacements = []

            for original, replacement in sorted_replacements:
                if original and original in new_full_text:
                    new_full_text = new_full_text.replace(original, replacement)
                    applied_replacements.append((original, replacement))
                    self.logger.info(f"  Replaced: '{original}' -> '{replacement}'")

            # If no exact matches, try fuzzy matching
            if not applied_replacements:
                fuzzy_results = self.text_processor.apply_fuzzy_replacements(
                    full_text, sorted_replacements
                )
                new_full_text = fuzzy_results["new_text"]
                applied_replacements = fuzzy_results["replacements"]

            # Update text frame if changes were made
            if applied_replacements and new_full_text != full_text:
                try:
                    # Store formatting from first run before clearing
                    original_font_props = {}
                    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                        first_run = text_frame.paragraphs[0].runs[0]
                        # Safe color handling
                        original_color = None
                        try:
                            if first_run.font.color and hasattr(first_run.font.color, 'rgb') and first_run.font.color.rgb:
                                original_color = first_run.font.color.rgb
                        except:
                            original_color = None
                            
                        original_font_props = {
                            'size': first_run.font.size,
                            'name': first_run.font.name,
                            'bold': first_run.font.bold,
                            'italic': first_run.font.italic,
                            'color': original_color
                        }

                    # Clear and rewrite the entire text frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = new_full_text

                    # Restore formatting if we captured it
                    if original_font_props and p.runs:
                        run = p.runs[0]
                        if original_font_props.get('size'):
                            run.font.size = original_font_props['size']
                        if original_font_props.get('name'):
                            run.font.name = original_font_props['name']
                        if original_font_props.get('bold') is not None:
                            run.font.bold = original_font_props['bold']
                        if original_font_props.get('italic') is not None:
                            run.font.italic = original_font_props['italic']
                        if original_font_props.get('color'):
                            try:
                                run.font.color.rgb = original_font_props['color']
                            except:
                                pass  # Skip color restoration if it fails

                    replacements_made = len(applied_replacements)
                    self.logger.info(
                        f"  Updated text frame: {replacements_made} replacements"
                    )
                    self.logger.debug(f"    Original: '{full_text}'")
                    self.logger.debug(f"    New: '{new_full_text}'")

                except Exception as e:
                    self.logger.error(f"Error updating text frame: {e}")
                    return 0

        return replacements_made

    def _apply_replacements_to_table(
        self, table, sorted_replacements: List[Tuple[str, str]]
    ) -> int:
        """
        Apply replacements to table cells.
        
        Iterates through all table rows and cells, applying text replacements
        to each cell's text frame while preserving table structure and formatting.
        
        Args:
            table: The python-pptx table object to modify
            sorted_replacements (List[Tuple[str, str]]): List of (original, replacement)
                tuples sorted by length (longest first)
                
        Returns:
            int: Total number of replacements made in all table cells
        """
        replacements_made = 0

        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    replacements_made += self._apply_replacements_to_text_frame(
                        cell.text_frame, sorted_replacements
                    )

        return replacements_made
